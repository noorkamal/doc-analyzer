# src/document_processor.py - Enhanced version for multiple files
import os
import tempfile
from typing import Dict, Any, List
import PyPDF2
import docx
from pptx import Presentation
import re
from datetime import datetime

class DocumentProcessor:
    def __init__(self, enable_privacy: bool = True):
        self.enable_privacy = enable_privacy
        self.supported_formats = ['pdf', 'docx', 'pptx']
    
    def process_multiple_documents(self, uploaded_files: List[Any], sanitization_level: str = "medium") -> List[Dict[str, Any]]:
        """Process multiple uploaded files"""
        processed_documents = []
        
        for uploaded_file in uploaded_files:
            try:
                doc_data = self.process_document(uploaded_file, sanitization_level)
                if doc_data:
                    doc_data['filename'] = uploaded_file.name
                    processed_documents.append(doc_data)
            except Exception as e:
                # Log error but continue with other documents
                print(f"Error processing {uploaded_file.name}: {str(e)}")
                continue
        
        return processed_documents
    
    def process_document(self, uploaded_file: Any, sanitization_level: str = "medium") -> Dict[str, Any]:
        """Process a single document file"""
        
        # Get file info
        file_name = uploaded_file.name
        file_extension = file_name.split('.')[-1].lower()
        file_size = len(uploaded_file.getvalue())
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Reset file pointer
        uploaded_file.seek(0)
        
        # Process based on file type
        document_data = {
            'filename': file_name,
            'type': file_extension,
            'file_size': file_size,
            'processed_at': datetime.now().isoformat(),
            'content': '',
            'metadata': {}
        }
        
        try:
            if file_extension == 'pdf':
                document_data.update(self._process_pdf(uploaded_file))
            elif file_extension == 'docx':
                document_data.update(self._process_docx(uploaded_file))
            elif file_extension == 'pptx':
                document_data.update(self._process_pptx(uploaded_file))
            
            # Apply privacy sanitization if enabled
            if self.enable_privacy and sanitization_level != "none":
                sanitized_data = self._sanitize_content(
                    document_data['content'], 
                    sanitization_level
                )
                document_data['content'] = sanitized_data['content']
                document_data['sanitization_info'] = sanitized_data['info']
            
            return document_data
            
        except Exception as e:
            raise Exception(f"Error processing {file_name}: {str(e)}")
    
    def _process_pdf(self, uploaded_file) -> Dict[str, Any]:
        """Extract content from PDF file"""
        content = ""
        pages = 0
        
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name
            
            # Read PDF
            with open(tmp_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                pages = len(pdf_reader.pages)
                
                for page_num in range(pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        content += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
            
            # Clean up
            os.unlink(tmp_file_path)
            
        except Exception as e:
            raise Exception(f"PDF processing error: {str(e)}")
        
        return {
            'content': content.strip(),
            'pages': pages,
            'metadata': {'total_pages': pages}
        }
    
    def _process_docx(self, uploaded_file) -> Dict[str, Any]:
        """Extract content from Word document"""
        content = ""
        paragraphs = []
        
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name
            
            # Read DOCX
            doc = docx.Document(tmp_file_path)
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
                    content += paragraph.text + "\n\n"
            
            # Clean up
            os.unlink(tmp_file_path)
            
        except Exception as e:
            raise Exception(f"DOCX processing error: {str(e)}")
        
        return {
            'content': content.strip(),
            'paragraphs': paragraphs,
            'metadata': {'total_paragraphs': len(paragraphs)}
        }
    
    def _process_pptx(self, uploaded_file) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation"""
        content = ""
        slide_count = 0
        slides_content = []
        
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name
            
            # Read PPTX
            prs = Presentation(tmp_file_path)
            slide_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    slides_content.append(slide_text)
                    content += slide_text + "\n"
            
            # Clean up
            os.unlink(tmp_file_path)
            
        except Exception as e:
            raise Exception(f"PPTX processing error: {str(e)}")
        
        return {
            'content': content.strip(),
            'slide_count': slide_count,
            'slides_content': slides_content,
            'metadata': {'total_slides': slide_count}
        }
    
    def _sanitize_content(self, content: str, level: str) -> Dict[str, Any]:
        """Sanitize content based on privacy level"""
        original_length = len(content)
        sanitized_content = content
        removed_items = {
            'emails': 0,
            'phone_numbers': 0,
            'credit_cards': 0,
            'ssns': 0,
            'names': 0,
            'addresses': 0
        }
        
        if level == "low":
            # Basic sanitization - only obvious patterns
            patterns = [
                (r'\b\d{4}[\s-]?\d{4}[\s-]?\d{4}[\s-]?\d{4}\b', '[CREDIT_CARD]', 'credit_cards'),
                (r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[PHONE]', 'phone_numbers'),
            ]
        
        elif level == "medium":
            # Standard sanitization
            patterns = [
                (r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL]', 'emails'),
                (r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[PHONE]', 'phone_numbers'),
                (r'\b\d{4}[\s-]?\d{4}[\s-]?\d{4}[\s-]?\d{4}\b', '[CREDIT_CARD]', 'credit_cards'),
            ]
        
        elif level == "high":
            # Aggressive sanitization
            patterns = [
                (r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL]', 'emails'),
                (r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[PHONE]', 'phone_numbers'),
                (r'\b\d{4}[\s-]?\d{4}[\s-]?\d{4}[\s-]?\d{4}\b', '[CREDIT_CARD]', 'credit_cards'),
                (r'\b\d{3}[-.]?\d{2}[-.]?\d{4}\b', '[SSN]', 'ssns'),
                (r'\b[A-Z][a-z]+ [A-Z][a-z]+\b', '[NAME]', 'names'),  # Simple name pattern
                (r'\b\d+\s+[A-Za-z\s]+(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Drive|Dr|Lane|Ln)\b', '[ADDRESS]', 'addresses'),
            ]
        else:
            patterns = []
        
        # Apply sanitization patterns
        for pattern, replacement, category in patterns:
            matches = re.findall(pattern, sanitized_content)
            removed_items[category] = len(matches)
            sanitized_content = re.sub(pattern, replacement, sanitized_content)
        
        return {
            'content': sanitized_content,
            'info': {
                'sanitization_level': level,
                'original_length': original_length,
                'sanitized_length': len(sanitized_content),
                'removed_items': removed_items,
                'processed_at': datetime.now().isoformat()
            }
        }
    
    def get_document_stats(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Get statistics across all processed documents"""
        if not documents:
            return {}
        
        total_docs = len(documents)
        total_words = sum(len(doc.get('content', '').split()) for doc in documents)
        file_types = {}
        total_size = 0
        
        for doc in documents:
            doc_type = doc.get('type', 'unknown')
            file_types[doc_type] = file_types.get(doc_type, 0) + 1
            total_size += doc.get('file_size', 0)
        
        return {
            'total_documents': total_docs,
            'total_words': total_words,
            'total_size_bytes': total_size,
            'file_types': file_types,
            'average_words_per_doc': total_words // total_docs if total_docs > 0 else 0
        }
